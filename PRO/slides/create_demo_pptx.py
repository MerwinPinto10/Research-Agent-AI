from pptx import Presentation

# Create a new presentation
prs = Presentation()

# Titles for the slides
titles = ["Problem", "Demo", "Tech", "Future"]

# Add each slide with a Title & Content layout
for title in titles:
    slide_layout = prs.slide_layouts[1]  # 0 = Title Slide, 1 = Title + Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Add content here"

# Save the presentation
prs.save("slides/demo.pptx")

print("demo.pptx created successfully!")
